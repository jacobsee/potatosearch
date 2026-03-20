"""
Plaintext / Markdown / general document file storage backend.

Recursively scans a directory for text and document files.
Serves as both a useful backend and a reference implementation for writing new plugins.

Supported formats:
  .txt
  .md
  .markdown
  .text
  .rst
  .pdf            – via pypdf
  .docx           – via python-docx
  .pptx           – via python-pptx
  .xlsx           – via openpyxl
  .odt .ods .odp  – via odfpy

Locator format: relative file (or directory) path from the configured root directory.:w
    e.g.  "notes/2024/quantum-computing.md"
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Iterator

from potatosearch.core import Document, StorageBackend

log = logging.getLogger(__name__)

_PLAIN_EXTENSIONS = {".txt", ".md", ".markdown", ".text", ".rst"}
_RICH_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".odt", ".ods", ".odp"}
_EXTENSIONS = _PLAIN_EXTENSIONS | _RICH_EXTENSIONS


# ---------------------------------------------------------------------------
# Per-format text extractors
# ---------------------------------------------------------------------------

def _read_plain(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader  # type: ignore[import]
    reader = PdfReader(path)
    parts = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p for p in parts if p.strip())


def _read_docx(path: Path) -> str:
    import docx  # type: ignore[import]
    doc = docx.Document(path)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def _read_pptx(path: Path) -> str:
    from pptx import Presentation  # type: ignore[import]
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _read_xlsx(path: Path) -> str:
    import openpyxl  # type: ignore[import]
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(v) for v in row if v is not None)
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


def _read_odf(path: Path) -> str:
    from odf import teletype  # type: ignore[import]
    from odf.opendocument import load  # type: ignore[import]
    doc = load(path)
    return teletype.extractText(doc.body)


_READERS = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".xlsx": _read_xlsx,
    ".odt": _read_odf,
    ".ods": _read_odf,
    ".odp": _read_odf,
}


def _read_file(path: Path) -> str:
    """Dispatch to the appropriate extractor based on file extension."""
    reader = _READERS.get(path.suffix.lower(), _read_plain)
    text = reader(path)
    # Sanitize lone surrogates that some extractors (especially pypdf) emit
    return text.encode("utf-8", errors="replace").decode("utf-8")


# ---------------------------------------------------------------------------
# Backend
# ---------------------------------------------------------------------------

class PlaintextBackend(StorageBackend):
    """
    Backend for directories of plain text, Markdown, and common document files.

    Args:
        root_dirs: Directories to scan recursively.
        extensions: File extensions to include (default: all supported formats).
    """

    def __init__(
        self,
        root_dirs: list[Path],
        extensions: set[str] | None = None,
        *,
        backend_id: str = "plaintext",
    ):
        self._backend_id = backend_id
        self._root_dirs = [Path(d) for d in root_dirs]
        self._extensions = extensions or _EXTENSIONS

    @property
    def name(self) -> str:
        return self._backend_id

    def iterate_documents(self) -> Iterator[Document]:
        for root in self._root_dirs:
            if not root.is_dir():
                log.warning("Skipping non-existent directory: %s", root)
                continue
            for path in sorted(root.rglob("*")):
                if path.suffix.lower() not in self._extensions:
                    continue
                if not path.is_file():
                    continue
                try:
                    text = _read_file(path)
                    if not text.strip():
                        continue
                    locator = str(path.relative_to(root))
                    yield Document(
                        locator=locator,
                        title=path.stem,
                        text=text,
                        metadata={"root": str(root), "extension": path.suffix},
                    )
                except ImportError as exc:
                    log.warning(
                        "Skipping %s - missing optional dependency: %s", path, exc
                    )
                except Exception:
                    log.warning("Failed to read %s", path, exc_info=True)

    def document_count_hint(self) -> int | None:
        total = 0
        for root in self._root_dirs:
            if not root.is_dir():
                continue
            for path in root.rglob("*"):
                if path.suffix.lower() in self._extensions and path.is_file():
                    total += 1
        return total or None

    def retrieve_text(self, locator: str, char_start: int, char_end: int) -> str:
        return self.retrieve_document(locator)[char_start:char_end]

    def retrieve_document(self, locator: str) -> str:
        for root in self._root_dirs:
            candidate = root / locator
            if candidate.is_file():
                return _read_file(candidate)
        raise FileNotFoundError(f"File not found in any root: {locator}")
